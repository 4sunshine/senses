import os

import numpy as np
from pptx import Presentation
from pptx.util import Pt
from omegaconf import OmegaConf
import textwrap as tw
import tempfile

from visio.video import VideoInput, AVStreamWriter, CV2WebCam, VideoDefaults, StaticInput, MPSimpleFaceDetector

import torch
import time
import cv2


class PPTToElements(object):
    DEFAULT_PT_SIZE = 2400

    def __init__(self, filepath, external_cfg):
        prs = Presentation(filepath)
        self.width, self.height = prs.slide_width, prs.slide_height
        self.target_w, self.target_h = external_cfg.size
        self.origin_x, self.origin_y = external_cfg.origin
        self.default_spacing = 1.5
        self.bullet = '\u2192'
        self.temp_mapping = dict()
        self.temp_dir = tempfile.TemporaryDirectory()
        self.all_texts = self.texts(prs)
        self.all_videos = self.videos(prs)
        print(self.all_texts)
        print(self.all_videos)
        self.close()

    @staticmethod
    def parse_alignment(alignment, default_align='l'):
        if alignment is None:
            alignment = default_align
        alignment = str(alignment).lower()
        if 'center' in alignment:
            alignment = 'ctr'
        return alignment

    def place_texts_in_shape(self, shape_bbox, texts_data, spacing, need_bullet):
        x_0, y_0, x_1, y_1 = shape_bbox
        w, h = x_1 - x_0, y_1 - y_0
        w = int(round(self.target_w * w / self.width))
        h = int(round(self.target_h * h / self.height))
        x_0 = int(round(self.target_w * x_0 / self.width))
        y_0 = int(round(self.target_h * y_0 / self.height))
        texts = texts_data['text']
        anchors = texts_data['anchor']
        aligns = texts_data['align']
        aligns = [self.parse_alignment(a) for a in aligns]
        font_sizes = texts_data['font_size']
        font_sizes = [int(round(self.target_h * f_s / self.height)) for f_s in font_sizes]
        n_letters = max([int(round(w / f_s)) for f_s in font_sizes])
        if need_bullet:
            texts = [tw.fill(' '.join([self.bullet, text]), n_letters) for text in texts]
        else:
            texts = [tw.fill(text, n_letters) for text in texts]
        target_height = sum([int(round(f_s * spacing * len(text.splitlines())))
                             for f_s, text in zip(font_sizes, texts)])
        if 'b' in anchors:
            texts_origin = y_0 + h - target_height
        else:
            texts_origin = y_0

        result_data = {
            'texts': texts,
            'aligns': aligns,
            'anchors': anchors,
            'origin': (x_0, texts_origin),
            'size': (w, target_height),
        }
        return result_data

    def close(self):
        self.temp_dir.cleanup()

    def videos(self, prs):
        all_videos = []
        for i, slide in enumerate(prs.slides):
            for k in slide.part.rels.keys():
                print(dir(k))
                print(k.rId)
                print(k.target_ref)
                if str(k.target_ref).endswith('.mp4'):
                    target_path = os.path.join(self.temp_dir.name, f'{k.rId}.mp4')
                    with open(target_path, 'wb') as f:
                        f.write(k.target_part.blob)
                    self.temp_mapping[str(k.rId)] = target_path

            slide_videos = []
            slide_bboxes = []
            for shape in slide.shapes:
                # print(shape.shape_type)  # TODO: CHECK LATER FOR A VIDEO/PIC SUPPORT
                shape_type = str(shape.shape_type).lower()

                if 'media' in shape_type:
                    x_0, y_0 = shape.left, shape.top
                    s_w, s_h = shape.width, shape.height
                    shape_bbox = [x_0, y_0, x_0 + s_w, y_0 + s_h]

                    for i in shape._pic.nvPicPr.nvPr.iterchildren():
                        if 'videoFile' in i.tag:
                            for k, v in i.items():
                                if 'link' in k:
                                    video_reference = v
                                    slide_videos.append(self.temp_mapping[str(v)])

                                    shape_bbox[0] = int(round(self.target_w * shape_bbox[0] / self.width))
                                    shape_bbox[2] = int(round(self.target_w * shape_bbox[2] / self.width))

                                    shape_bbox[1] = int(round(self.target_h * shape_bbox[1] / self.height))
                                    shape_bbox[3] = int(round(self.target_h * shape_bbox[3] / self.height))

                                    slide_bboxes.append(shape_bbox)

            video_data = {
                'videos': slide_videos,
                'bboxes': slide_bboxes,
            }
            all_videos.append(video_data)

        return all_videos


    def texts(self, prs):
        def get_level_parameters():
            pass

        def get_defaults(shapes, placeholders):
            text_shapes_ids = [shape.shape_id for shape in shapes if shape.has_text_frame]
            slide_defaults = dict()
            for p in placeholders:
                if p.shape_id not in text_shapes_ids:
                    continue
                sizes = []
                aligns = []
                anchors = []
                txBody = p.element.txBody
                anchor = txBody.bodyPr.get('anchor')
                for child in txBody.iterchildren():
                    if 'lstStyle' in child.tag:
                        for style in child.iterchildren():
                            align = style.get('algn')
                            for props in style.iterchildren():
                                if 'defRPr' in props.tag:
                                    def_font_size = Pt(int(props.get('sz')) / 100)
                                    cur_align = props.get('algn')
                                    if cur_align is None:
                                        aligns.append(align)
                                    else:
                                        aligns.append(cur_align)
                                    if def_font_size is None:
                                        sizes.append(self.height // 10)
                                    else:
                                        sizes.append(def_font_size)
                                    anchors.append(anchor)  # TODO: CHECK THIS

                slide_defaults[p.shape_id] = {
                    'anchors': anchors,
                    'aligns': aligns,
                    'sizes': sizes,
                }
            return slide_defaults

        def get_master_defaults(master_slide):
            # TODO: LATER IMPLEMENT MAPPING
            # key_mapping = {'title': 2, 'body': 3}
            # print(dir(master_slide))
            # print(master_slide.shapes._spTree.xml)
            # print(dir(master_slide.shapes._spTree))
            # for s in master_slide.shapes._spTree.iterchildren():
            #     if '}sp' in s.tag:
            #         for s_child in s.iterchildren():
            #             if 'nvSpPr' in s_child.tag:
            #                 #print(dir(s_child))
            #                 print(dir(s_child))
            # raise
            #         # print(dir(s))
            #         # print(dir(s.nvSpPr))
            #         # print(s.nvSpPr.shape_id)
            #         # print(s.nvSpPr.shape_name)
            # # raise
            master_defaults = dict()
            for p in master_slide.element.iterchildren():
                if 'txStyles' in p.tag:
                    for child in p.iterchildren():  ## ITERATE OVER (TITLE STYLE, BODY, OTHER)
                        typename = child.tag.split('}')[-1].replace('Style', '')
                        aligns = []
                        sizes = []
                        anchors = []
                        for style in child.iterchildren():  ## ITERATE OVER LEVELS
                            align = style.get('algn')
                            for props in style.iterchildren():
                                if 'defRPr' in props.tag:# and ('lvl' in props.tag):
                                    size = props.get('sz')
                                    if size is None:
                                        size = self.DEFAULT_PT_SIZE
                                    def_font_size = Pt(int(size) / 100)
                                    cur_align = props.get('algn')
                                    if cur_align is None:
                                        aligns.append(align)
                                    else:
                                        aligns.append(cur_align)
                                    if def_font_size is None:
                                        sizes.append(self.height // 10)
                                    else:
                                        sizes.append(def_font_size)
                                    anchors.append(None)  # TODO: CHECK THIS

                        master_defaults[typename] = {
                            'anchors': anchors,
                            'aligns': aligns,
                            'sizes': sizes,
                        }

            def mapping(id):
                if id == 2:
                    return 'title'
                elif id == 3:
                    return 'body'
                else:
                    return 'other'

            return master_defaults, mapping

        all_texts = []
        master_defaults, master_mapping = get_master_defaults(prs.slide_master)

        for i, slide in enumerate(prs.slides):
            defaults = get_defaults(slide.shapes, slide.slide_layout.placeholders)

            slide_texts = []
            for shape in slide.shapes:
                # print(shape.shape_type)  # TODO: CHECK LATER FOR A VIDEO/PIC SUPPORT
                if not shape.has_text_frame:
                    continue
                #print(shape.shape_id)
                # print(dir(shape))
                # print(dir(shape._sp))
                # print(shape._sp._get_xfrm_attr('type'))

                x_0, y_0 = shape.left, shape.top
                s_w, s_h = shape.width, shape.height

                shape_bbox = [x_0, y_0, x_0 + s_w, y_0 + s_h]

                need_bullet = len(shape.text_frame.paragraphs) > 1

                shape_texts = []
                anchors = []
                aligns = []
                font_sizes = []

                for paragraph in shape.text_frame.paragraphs:
                    spacing = paragraph.line_spacing
                    if spacing is None:
                        spacing = self.default_spacing  # TODO: CHECK DEFAULTS
                    align = paragraph.alignment

                    for run in paragraph.runs:
                        text = run.text
                        font_size = run.font.size
                        # print(text)
                        # print(dir(run))
                        anchor = None
                        if font_size is None:
                            if len(defaults[shape.shape_id]['sizes']) > 0:
                                font_size = defaults[shape.shape_id]['sizes'][paragraph.level]
                                align = defaults[shape.shape_id]['aligns'][paragraph.level]
                                anchor = defaults[shape.shape_id]['anchors'][paragraph.level]
                            else:
                                font_size = master_defaults[master_mapping(shape.shape_id)]['sizes'][paragraph.level]
                                align = master_defaults[master_mapping(shape.shape_id)]['aligns'][paragraph.level]
                                anchor = master_defaults[master_mapping(shape.shape_id)]['anchors'][paragraph.level]

                        shape_texts.append(text)
                        anchors.append(anchor)
                        aligns.append(align)
                        font_sizes.append(font_size)

                shape_texts_data = {
                    'text': shape_texts,
                    'anchor': anchors,
                    'align': aligns,
                    'font_size': font_sizes,
                }

                formatted_shape_texts_data = self.place_texts_in_shape(shape_bbox, shape_texts_data,
                                                                       spacing, need_bullet)  # TODO: FIX SPACING

                slide_texts.append(formatted_shape_texts_data)

            all_texts.append(slide_texts)

        return all_texts


def test_ppt_class(path):
    from visio.video import VideoDefaults
    cfg = VideoDefaults()
    s = PPTToElements(path, cfg)

